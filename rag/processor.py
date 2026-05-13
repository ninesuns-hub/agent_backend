import os
import logging
from typing import List, Dict, Any
from pptx import Presentation
from agent_core.config.settings import settings

logger = logging.getLogger(__name__)

class PPTXParser:
    """
    PPTX 解析器：将 PPTX 转换为统一的 Markdown 格式
    """
    def __init__(self, output_dir: str = None):
        self.output_dir = output_dir or settings.PROCESSED_DATA_DIR
        os.makedirs(self.output_dir, exist_ok=True)

    def parse(self, file_path: str) -> List[Dict[str, Any]]:
        """
        解析单个 PPTX 文件
        返回: List of { "text": str, "page": int, "metadata": dict }
        """
        if not os.path.exists(file_path):
            logger.error(f"文件不存在: {file_path}")
            return []

        prs = Presentation(file_path)
        file_name = os.path.basename(file_path)
        chunks = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            content = "\n".join(slide_text)
            if content:
                # 构造统一的 Markdown 格式字符串
                md_content = f"## Slide {i+1} from {file_name}\n\n{content}"
                
                chunks.append({
                    "text": md_content,
                    "page": i + 1,
                    "source_file": file_name,
                    "source_type": "pptx",
                    "metadata": {
                        "total_slides": len(prs.slides)
                    }
                })
        
        # 保存为本地 Markdown 文件以便调试
        self._save_to_markdown(file_name, chunks)
        return chunks

    def _save_to_markdown(self, file_name: str, chunks: List[Dict[str, Any]]):
        md_file_path = os.path.join(self.output_dir, f"{file_name}.md")
        try:
            with open(md_file_path, "w", encoding="utf-8") as f:
                for chunk in chunks:
                    f.write(chunk["text"])
                    f.write("\n\n---\n\n")
            logger.info(f"已保存处理后的 Markdown 文件: {md_file_path}")
        except Exception as e:
            logger.error(f"保存 Markdown 失败: {e}")

class DocumentProcessor:
    """
    统一文档处理器：处理不同格式的文件并入库
    """
    def __init__(self):
        self.pptx_parser = PPTXParser()
        # 未来可以在这里增加 PDFParser, DocxParser 等

    def process_directory(self, directory: str) -> List[Dict[str, Any]]:
        all_chunks = []
        if not os.path.exists(directory):
            logger.warning(f"目录不存在: {directory}")
            return []

        for root, _, files in os.walk(directory):
            for file in files:
                file_path = os.path.join(root, file)
                if file.endswith(".pptx"):
                    logger.info(f"正在处理 PPTX: {file}")
                    chunks = self.pptx_parser.parse(file_path)
                    all_chunks.extend(chunks)
                # 未来增加其他格式的判断
        
        return all_chunks

if __name__ == "__main__":
    # 简单的测试逻辑
    logging.basicConfig(level=logging.INFO)
    processor = DocumentProcessor()
    # 假设 backend/data/pptx 目录下有文件
    # chunks = processor.process_directory(settings.RAW_DATA_DIR)
    # print(f"Total chunks: {len(chunks)}")
